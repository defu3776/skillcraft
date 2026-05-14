"""Portable Office deep parsing runtime for cross-repo skill sharing."""

from __future__ import annotations

import csv
import hashlib
import json
import re
import shutil
import subprocess
from dataclasses import asdict, dataclass, field
from datetime import datetime
from pathlib import Path
from tempfile import TemporaryDirectory
from typing import Any

from openpyxl import load_workbook

from runtime.executables import find_executable


SUBPROCESS_TIMEOUT_SECONDS = 120
MAX_PDF_OCR_PAGES = 25


@dataclass(slots=True)
class PipelineConfig:
    input_path: Path
    output_root: Path
    enable_markitdown: bool = True
    enable_visual_export: bool = True
    enable_ocr: bool = True
    ocr_backend: str = "local"
    recurse: bool = True
    extract_attachments: bool = True
    supported_office_exts: tuple[str, ...] = (
        ".xlsx",
        ".xlsm",
        ".xls",
        ".csv",
        ".docx",
        ".doc",
        ".pptx",
        ".ppt",
    )
    attachment_exts: tuple[str, ...] = (".pdf", ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp")
    archive_exts: tuple[str, ...] = (".zip", ".7z", ".rar")
    skip_prefixes: tuple[str, ...] = ("~$", ".~")
    skip_names: tuple[str, ...] = (".ds_store", "thumbs.db", "desktop.ini")
    validation_warnings: list[str] = field(default_factory=list)

    def normalize(self) -> "PipelineConfig":
        self.validation_warnings.clear()
        self.input_path = self.input_path.expanduser().resolve()
        self.output_root = self.output_root.expanduser().resolve()
        if not self.input_path.exists():
            raise FileNotFoundError(f"input path does not exist: {_path_for_message(self.input_path)}")
        if not (self.input_path.is_file() or self.input_path.is_dir()):
            raise ValueError(f"input path must be a file or directory: {_path_for_message(self.input_path)}")
        if self.input_path.is_file() and self.output_root == self.input_path:
            raise ValueError("output root must be a directory, not the input file")
        if self.input_path.is_dir() and self.output_root == self.input_path:
            raise ValueError("output root must not be the same directory as input path")
        if self.input_path.is_dir() and _is_relative_to(self.output_root, self.input_path):
            self.validation_warnings.append(
                "output_root is inside input_path; generated output files are excluded from input inventory"
            )
        return self

    def ensure_output_dirs(self) -> None:
        for path in (
            self.output_root,
            self.extracted_markdown_path,
            self.visual_exports_path,
            self.ocr_results_path,
            self.deep_notes_path,
        ):
            path.mkdir(parents=True, exist_ok=True)

    @property
    def extracted_markdown_path(self) -> Path:
        return self.output_root / "extracted_markdown"

    @property
    def visual_exports_path(self) -> Path:
        return self.output_root / "visual_exports"

    @property
    def ocr_results_path(self) -> Path:
        return self.output_root / "ocr_results"

    @property
    def deep_notes_path(self) -> Path:
        return self.output_root / "deep_reading_notes"


@dataclass(slots=True)
class FileInventoryRow:
    relative_path: str
    file_name: str
    file_type: str
    size_bytes: int
    modified_time: str
    should_process: bool
    skip_reason: str = ""
    tags: list[str] = field(default_factory=list)


@dataclass(slots=True)
class MarkdownExtractionResult:
    source_file: str
    markdown_path: str | None
    status: str
    log: str


@dataclass(slots=True)
class CellRecord:
    coordinate: str
    value: str | None
    data_type: str | None
    formula: str | None = None
    cached_value: str | None = None
    is_merged_anchor: bool = False
    merged_range: str | None = None
    comment: str | None = None
    hyperlink: str | None = None


@dataclass(slots=True)
class ObjectRecord:
    object_type: str
    sheet_name: str
    anchor: str | None = None
    width: int | None = None
    height: int | None = None
    description: str | None = None
    export_path: str | None = None


@dataclass(slots=True)
class SheetAnalysis:
    sheet_name: str
    index: int
    state: str
    dimensions: str
    max_row: int
    max_column: int
    freeze_panes: str | None
    print_area: str | None
    print_titles: str | None = None
    page_setup: dict[str, str] = field(default_factory=dict)
    merged_ranges: list[str] = field(default_factory=list)
    table_names: list[str] = field(default_factory=list)
    named_range_refs: list[str] = field(default_factory=list)
    data_validation_count: int = 0
    conditional_format_count: int = 0
    chart_count: int = 0
    image_count: int = 0
    drawing_count: int = 0
    embedded_object_count: int = 0
    cell_records: list[CellRecord] = field(default_factory=list)
    object_records: list[ObjectRecord] = field(default_factory=list)


@dataclass(slots=True)
class WorkbookAnalysis:
    source_file: str
    workbook_name: str
    sheet_order: list[str]
    hidden_sheets: list[str]
    named_ranges: dict[str, str]
    sheet_analyses: list[SheetAnalysis]
    extraction_warnings: list[str] = field(default_factory=list)


@dataclass(slots=True)
class DocumentSectionRecord:
    section_type: str
    location: str
    text: str


@dataclass(slots=True)
class DocumentAnalysis:
    source_file: str
    document_name: str
    document_type: str
    sections: list[DocumentSectionRecord] = field(default_factory=list)
    image_count: int = 0
    extraction_warnings: list[str] = field(default_factory=list)


@dataclass(slots=True)
class OCRResult:
    source_visual_path: str
    status: str
    ocr_text: str
    model_notes: str = ""
    page_ref: str | None = None


@dataclass(slots=True)
class PipelineResult:
    created_at: str
    input_path: str
    output_root: str
    file_inventory: list[FileInventoryRow] = field(default_factory=list)
    markdown_results: list[MarkdownExtractionResult] = field(default_factory=list)
    workbook_results: list[WorkbookAnalysis] = field(default_factory=list)
    document_results: list[DocumentAnalysis] = field(default_factory=list)
    ocr_results: list[OCRResult] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)

    @classmethod
    def bootstrap(cls, input_path: Path, output_root: Path) -> "PipelineResult":
        return cls(
            created_at=datetime.utcnow().isoformat(timespec="seconds") + "Z",
            input_path=input_path.name or str(input_path),
            output_root=output_root.name or "output_root",
        )


def _file_type(path: Path) -> str:
    return path.suffix.lower().lstrip(".") or "unknown"


def _path_for_message(path: Path) -> str:
    return path.name or str(path)


def _is_relative_to(path: Path, base: Path) -> bool:
    try:
        path.relative_to(base)
        return True
    except ValueError:
        return False


def _safe_token(value: str, max_length: int = 80) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9._-]+", "_", value.strip())
    cleaned = cleaned.strip("._-") or "item"
    return cleaned[:max_length].rstrip("._-") or "item"


def _artifact_stem(relative_path: str) -> str:
    normalized = relative_path.replace("\\", "/")
    parts = [part for part in normalized.split("/") if part and part not in {".", ".."}]
    base = "__".join(_safe_token(part) for part in parts) or "artifact"
    digest = hashlib.sha256(normalized.encode("utf-8")).hexdigest()[:8]
    return f"{base}__{digest}"


def _md_cell(value: object) -> str:
    text = "" if value is None else str(value)
    return text.replace("\n", " ").replace("|", "\\|")


def _display_output_path(path: Path, output_root: Path) -> str:
    resolved_path = path.expanduser().resolve()
    resolved_root = output_root.expanduser().resolve()
    if _is_relative_to(resolved_path, resolved_root):
        return str(resolved_path.relative_to(resolved_root))
    return path.name


def _run_tool(args: list[str], timeout_seconds: int = SUBPROCESS_TIMEOUT_SECONDS) -> tuple[int, str]:
    try:
        completed = subprocess.run(
            args,
            capture_output=True,
            text=True,
            check=False,
            timeout=timeout_seconds,
        )
    except subprocess.TimeoutExpired:
        return 124, f"{Path(args[0]).name} timed out after {timeout_seconds}s"
    log = (completed.stdout or "") + ("\n" + completed.stderr if completed.stderr else "")
    return completed.returncode, log.strip()


def _redact_paths(log: str, *paths: Path) -> str:
    redacted = log
    for path in paths:
        resolved = str(path.expanduser().resolve())
        redacted = redacted.replace(resolved, path.name)
    return redacted


def _convert_with_soffice(source: Path, target_ext: str, work_dir: Path) -> Path | None:
    soffice = find_executable("soffice")
    if not soffice:
        return None
    returncode, _log = _run_tool(
        [soffice, "--headless", "--convert-to", target_ext.lstrip("."), "--outdir", str(work_dir), str(source)],
    )
    if returncode != 0:
        return None
    candidate = work_dir / f"{source.stem}{target_ext}"
    return candidate if candidate.exists() else None


def scan_input(config: PipelineConfig) -> list[FileInventoryRow]:
    root = config.input_path
    if root.is_file():
        roots = [root]
        base = root.parent
    else:
        roots = list(root.rglob("*")) if config.recurse else list(root.glob("*"))
        roots = [p for p in roots if p.is_file() and not _is_relative_to(p.resolve(), config.output_root)]
        base = root
    rows: list[FileInventoryRow] = []
    for path in roots:
        lower_name = path.name.lower()
        ext = path.suffix.lower()
        should_process = False
        skip_reason = ""
        tags: list[str] = []
        if lower_name in config.skip_names or any(path.name.startswith(prefix) for prefix in config.skip_prefixes):
            skip_reason = "temporary_or_system_file"
        elif ext in config.supported_office_exts:
            should_process = True
            tags.append("office_target")
            if ext in {".xlsx", ".xlsm", ".xls", ".csv"}:
                tags.append("spreadsheet_target")
            if ext in {".doc", ".docx", ".ppt", ".pptx"}:
                tags.append("document_target")
        elif ext in config.attachment_exts:
            should_process = config.extract_attachments
            tags.append("attachment_candidate")
            if not should_process:
                skip_reason = "attachment_disabled"
        elif ext in config.archive_exts:
            skip_reason = "archive_not_expanded_pending_confirmation"
            tags.append("archive")
        else:
            skip_reason = "unsupported_type"
        stat = path.stat()
        rows.append(
            FileInventoryRow(
                relative_path=str(path.relative_to(base)),
                file_name=path.name,
                file_type=_file_type(path),
                size_bytes=stat.st_size,
                modified_time=datetime.fromtimestamp(stat.st_mtime).isoformat(timespec="seconds"),
                should_process=should_process,
                skip_reason=skip_reason,
                tags=tags,
            )
        )
    rows.sort(key=lambda x: x.relative_path.lower())
    return rows


def filter_spreadsheet_targets(rows: list[FileInventoryRow]) -> list[FileInventoryRow]:
    return [row for row in rows if row.should_process and "spreadsheet_target" in row.tags]


def filter_document_targets(rows: list[FileInventoryRow]) -> list[FileInventoryRow]:
    return [row for row in rows if row.should_process and "document_target" in row.tags]


def filter_attachment_targets(rows: list[FileInventoryRow]) -> list[FileInventoryRow]:
    return [row for row in rows if row.should_process and "attachment_candidate" in row.tags]


def extract_to_markdown(
    source_file: Path,
    output_file: Path,
    source_label: str | None = None,
    output_label: str | None = None,
) -> MarkdownExtractionResult:
    cli = shutil.which("markitdown")
    if cli is None:
        return MarkdownExtractionResult(
            source_file=source_label or source_file.name,
            markdown_path=None,
            status="skipped",
            log="markitdown CLI not found in PATH",
        )
    output_file.parent.mkdir(parents=True, exist_ok=True)
    returncode, log = _run_tool(
        [cli, str(source_file), "-o", str(output_file)],
    )
    log = _redact_paths(log, source_file, output_file)
    if returncode == 0 and output_file.exists():
        return MarkdownExtractionResult(
            source_label or source_file.name,
            output_label or output_file.name,
            "success",
            log.strip(),
        )
    return MarkdownExtractionResult(
        source_label or source_file.name,
        None,
        "failed",
        (log or "markitdown failed").strip(),
    )


def _to_text(value: Any) -> str | None:
    return None if value is None else str(value)


def _column_name(index: int) -> str:
    name = ""
    while index:
        index, rem = divmod(index - 1, 26)
        name = chr(65 + rem) + name
    return name or "A"


def _chart_title(chart: Any) -> str | None:
    title = getattr(chart, "title", None)
    if title is None:
        return None
    try:
        parts: list[str] = []
        for paragraph in title.tx.rich.p:
            for run in getattr(paragraph, "r", []):
                if getattr(run, "t", None):
                    parts.append(run.t)
        return " ".join(parts).strip() or None
    except Exception:
        return str(title)


def _print_titles(sheet: Any) -> str | None:
    parts = []
    if sheet.print_title_rows:
        parts.append(str(sheet.print_title_rows))
    if sheet.print_title_cols:
        parts.append(str(sheet.print_title_cols))
    return ", ".join(parts) if parts else None


def parse_csv_as_sheet(csv_path: Path) -> SheetAnalysis:
    rows: list[list[str]] = []
    with csv_path.open("r", encoding="utf-8-sig", newline="") as f:
        rows = list(csv.reader(f))
    max_row = len(rows)
    max_col = max((len(r) for r in rows), default=0)
    cells: list[CellRecord] = []
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row, start=1):
            if value:
                cells.append(CellRecord(f"{_column_name(col_idx)}{row_idx}", value, "s"))
    return SheetAnalysis(
        sheet_name="CSV",
        index=0,
        state="visible",
        dimensions=f"A1:{_column_name(max_col)}{max_row}" if max_row and max_col else "A1:A1",
        max_row=max_row,
        max_column=max_col,
        freeze_panes=None,
        print_area=None,
        cell_records=cells,
    )


def _analyze_xlsx_like(path: Path, source_label: str) -> WorkbookAnalysis:
    try:
        wb = load_workbook(filename=path, data_only=False, read_only=False, keep_links=False)
        wb_data = load_workbook(filename=path, data_only=True, read_only=False, keep_links=False)
    except Exception as exc:
        return WorkbookAnalysis(source_label, Path(source_label).name, [], [], {}, [], [f"workbook load failed: {exc}"])

    named_ranges: dict[str, str] = {}
    try:
        for name in wb.defined_names.definedName:
            named_ranges[name.name] = name.attr_text or ""
    except Exception:
        pass

    sheet_analyses: list[SheetAnalysis] = []
    for idx, sheet in enumerate(wb.worksheets):
        data_sheet = wb_data[sheet.title] if sheet.title in wb_data.sheetnames else None
        anchor_map: dict[str, str] = {}
        for merged in sheet.merged_cells.ranges:
            anchor = merged.start_cell.coordinate
            for row in range(merged.min_row, merged.max_row + 1):
                for col in range(merged.min_col, merged.max_col + 1):
                    anchor_map[sheet.cell(row=row, column=col).coordinate] = anchor
        cell_records: list[CellRecord] = []
        for row in sheet.iter_rows(min_row=1, max_row=sheet.max_row, min_col=1, max_col=sheet.max_column):
            for cell in row:
                if cell.value is None and cell.comment is None and cell.hyperlink is None:
                    continue
                cached = data_sheet[cell.coordinate].value if data_sheet is not None else None
                merge_anchor = anchor_map.get(cell.coordinate)
                merged_range = None
                is_anchor = False
                if merge_anchor:
                    is_anchor = merge_anchor == cell.coordinate
                    if is_anchor:
                        for merged in sheet.merged_cells.ranges:
                            if merged.start_cell.coordinate == merge_anchor:
                                merged_range = str(merged)
                                break
                cell_records.append(
                    CellRecord(
                        coordinate=cell.coordinate,
                        value=_to_text(cell.value),
                        data_type=cell.data_type,
                        formula=cell.value if cell.data_type == "f" else None,
                        cached_value=_to_text(cached),
                        is_merged_anchor=is_anchor,
                        merged_range=merged_range,
                        comment=cell.comment.text if cell.comment else None,
                        hyperlink=cell.hyperlink.target if cell.hyperlink and cell.hyperlink.target else None,
                    )
                )
        object_records: list[ObjectRecord] = []
        for chart in getattr(sheet, "_charts", []):
            anchor = None
            if hasattr(chart, "anchor") and getattr(chart.anchor, "_from", None):
                marker = chart.anchor._from
                anchor = f"{marker.col + 1},{marker.row + 1}"
            object_records.append(ObjectRecord("chart", sheet.title, anchor=anchor, description=_chart_title(chart)))
        for image in getattr(sheet, "_images", []):
            anchor = None
            if hasattr(image, "anchor") and getattr(image.anchor, "_from", None):
                marker = image.anchor._from
                anchor = f"{marker.col + 1},{marker.row + 1}"
            object_records.append(
                ObjectRecord(
                    "image",
                    sheet.title,
                    anchor=anchor,
                    width=getattr(image, "width", None),
                    height=getattr(image, "height", None),
                )
            )
        refs = [f"{k}: {v}" for k, v in named_ranges.items() if sheet.title in v]
        sheet_analyses.append(
            SheetAnalysis(
                sheet_name=sheet.title,
                index=idx,
                state=sheet.sheet_state,
                dimensions=sheet.calculate_dimension(),
                max_row=sheet.max_row,
                max_column=sheet.max_column,
                freeze_panes=str(sheet.freeze_panes) if sheet.freeze_panes else None,
                print_area=str(sheet.print_area) if sheet.print_area else None,
                print_titles=_print_titles(sheet),
                page_setup={
                    "orientation": str(sheet.page_setup.orientation) if sheet.page_setup else "",
                    "paper_size": str(sheet.page_setup.paperSize) if sheet.page_setup else "",
                },
                merged_ranges=[str(rng) for rng in sheet.merged_cells.ranges],
                table_names=list(sheet.tables.keys()),
                named_range_refs=refs,
                data_validation_count=len(sheet.data_validations.dataValidation),
                conditional_format_count=len(sheet.conditional_formatting),
                chart_count=len(getattr(sheet, "_charts", [])),
                image_count=len(getattr(sheet, "_images", [])),
                drawing_count=1 if getattr(sheet, "_drawing", None) else 0,
                embedded_object_count=0,
                cell_records=cell_records,
                object_records=object_records,
            )
        )
    hidden = [sheet.title for sheet in wb.worksheets if sheet.sheet_state != "visible"]
    return WorkbookAnalysis(source_label, Path(source_label).name, [s.title for s in wb.worksheets], hidden, named_ranges, sheet_analyses, [])


def analyze_workbook(path: Path, source_label: str | None = None) -> WorkbookAnalysis:
    label = source_label or path.name
    suffix = path.suffix.lower()
    if suffix == ".csv":
        sheet = parse_csv_as_sheet(path)
        return WorkbookAnalysis(label, Path(label).name, [sheet.sheet_name], [], {}, [sheet], [])
    if suffix in {".xlsx", ".xlsm"}:
        return _analyze_xlsx_like(path, label)
    if suffix == ".xls":
        with TemporaryDirectory(prefix="office_skill_xls_") as temp_dir:
            converted = _convert_with_soffice(path, ".xlsx", Path(temp_dir))
            if converted is None:
                return WorkbookAnalysis(
                    label,
                    Path(label).name,
                    [],
                    [],
                    {},
                    [],
                    [".xls conversion to .xlsx failed; install/configure LibreOffice soffice."],
                )
            analyzed = _analyze_xlsx_like(converted, label)
            analyzed.extraction_warnings.append("source .xls converted to .xlsx for deep parse")
            return analyzed
    return WorkbookAnalysis(label, Path(label).name, [], [], {}, [], [f"unsupported spreadsheet extension: {suffix}"])


def analyze_document(path: Path, source_label: str | None = None) -> DocumentAnalysis:
    suffix = path.suffix.lower()
    label = source_label or path.name
    analysis = DocumentAnalysis(source_file=label, document_name=Path(label).name, document_type=suffix.lstrip("."))
    if suffix == ".docx":
        try:
            import docx  # type: ignore

            doc = docx.Document(str(path))
            for idx, para in enumerate(doc.paragraphs, start=1):
                text = para.text.strip()
                if text:
                    analysis.sections.append(DocumentSectionRecord("paragraph", f"p{idx}", text))
            for t_idx, table in enumerate(doc.tables, start=1):
                for r_idx, row in enumerate(table.rows, start=1):
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text and cell.text.strip())
                    if row_text:
                        analysis.sections.append(DocumentSectionRecord("table_row", f"t{t_idx}r{r_idx}", row_text))
            rels = doc.part.rels
            analysis.image_count = sum(1 for rel in rels.values() if "image" in rel.reltype)
        except Exception as exc:
            analysis.extraction_warnings.append(f"docx parse failed: {exc}")
        return analysis
    if suffix == ".pptx":
        try:
            from pptx import Presentation  # type: ignore

            prs = Presentation(str(path))
            for s_idx, slide in enumerate(prs.slides, start=1):
                for sh_idx, shape in enumerate(slide.shapes, start=1):
                    text = ""
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                    if text:
                        analysis.sections.append(DocumentSectionRecord("slide_text", f"s{s_idx}sh{sh_idx}", text))
                    if getattr(shape, "shape_type", None) == 13:
                        analysis.image_count += 1
                if slide.has_notes_slide:
                    note_text = slide.notes_slide.notes_text_frame.text.strip()
                    if note_text:
                        analysis.sections.append(DocumentSectionRecord("slide_note", f"s{s_idx}note", note_text))
        except Exception as exc:
            analysis.extraction_warnings.append(f"pptx parse failed: {exc}")
        return analysis
    if suffix in {".doc", ".ppt"}:
        target_ext = ".docx" if suffix == ".doc" else ".pptx"
        with TemporaryDirectory(prefix="office_skill_doc_") as temp_dir:
            converted = _convert_with_soffice(path, target_ext, Path(temp_dir))
            if converted is None:
                analysis.extraction_warnings.append(f"{suffix} conversion to {target_ext} failed; install/configure soffice.")
                return analysis
            converted_analysis = analyze_document(converted, label)
            converted_analysis.source_file = label
            converted_analysis.document_name = Path(label).name
            converted_analysis.extraction_warnings.append(f"source {suffix} converted to {target_ext} for deep parse")
            return converted_analysis
    analysis.extraction_warnings.append(f"unsupported document extension: {suffix}")
    return analysis


def export_visual_assets(source_file: Path, workbook_analysis: WorkbookAnalysis, output_dir: Path) -> list[ObjectRecord]:
    exported: list[ObjectRecord] = []
    output_dir.mkdir(parents=True, exist_ok=True)
    if source_file.suffix.lower() in {".xlsx", ".xlsm"}:
        wb = load_workbook(filename=source_file, data_only=False, read_only=False, keep_links=False)
        for sheet in wb.worksheets:
            for idx, image in enumerate(getattr(sheet, "_images", []), start=1):
                image_path = output_dir / f"{_safe_token(source_file.stem)}__{_safe_token(sheet.title)}__img{idx}.png"
                try:
                    image_path.write_bytes(image._data())
                    exported.append(
                        ObjectRecord(
                            object_type="image_export",
                            sheet_name=sheet.title,
                            description="Embedded image extracted from workbook",
                            export_path=str(image_path),
                        )
                    )
                except Exception as exc:
                    workbook_analysis.extraction_warnings.append(f"image export failed on {sheet.title}#{idx}: {exc}")
    soffice = find_executable("soffice")
    if soffice:
        returncode, _log = _run_tool(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(output_dir), str(source_file)],
        )
        candidate = output_dir / f"{source_file.stem}.pdf"
        if returncode == 0 and candidate.exists():
            exported.append(
                ObjectRecord(
                    object_type="sheet_pdf_export",
                    sheet_name="*",
                    description="Workbook visual export via LibreOffice",
                    export_path=str(candidate),
                )
            )
    return exported


def export_document_visuals(source_file: Path, output_dir: Path) -> list[Path]:
    output_dir.mkdir(parents=True, exist_ok=True)
    exported: list[Path] = []
    soffice = find_executable("soffice")
    if not soffice:
        return exported
    returncode, _log = _run_tool(
        [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(output_dir), str(source_file)],
    )
    candidate = output_dir / f"{source_file.stem}.pdf"
    if returncode == 0 and candidate.exists():
        exported.append(candidate)
    return exported


def _run_local_tesseract(image_path: Path) -> OCRResult:
    try:
        import pytesseract
        from PIL import Image
    except Exception as exc:
        return OCRResult(str(image_path), "skipped", "", f"local OCR unavailable: {exc}")
    try:
        text = pytesseract.image_to_string(Image.open(image_path))
        return OCRResult(str(image_path), "success", text.strip(), "tesseract")
    except Exception as exc:
        return OCRResult(str(image_path), "failed", "", f"OCR failed: {exc}")


def _run_pdf_ocr(pdf_path: Path, backend: str) -> list[OCRResult]:
    if backend != "local":
        return [OCRResult(str(pdf_path), "skipped", "", f"unsupported backend: {backend}")]
    try:
        import pypdfium2 as pdfium
    except Exception as exc:
        return [OCRResult(str(pdf_path), "skipped", "", f"pdf OCR unavailable: {exc}")]
    results: list[OCRResult] = []
    pdf = pdfium.PdfDocument(str(pdf_path))
    page_count = len(pdf)
    for index in range(min(page_count, MAX_PDF_OCR_PAGES)):
        page = pdf[index]
        tmp_image = pdf_path.parent / f"{pdf_path.stem}.page{index + 1}.png"
        try:
            page.render(scale=2.0).to_pil().save(tmp_image)
            result = _run_local_tesseract(tmp_image)
            result.source_visual_path = str(pdf_path)
            result.page_ref = f"page_{index + 1}"
            results.append(result)
        except Exception as exc:
            results.append(
                OCRResult(
                    source_visual_path=str(pdf_path),
                    status="failed",
                    ocr_text="",
                    model_notes=f"pdf page OCR failed: {exc}",
                    page_ref=f"page_{index + 1}",
                )
            )
        finally:
            try:
                if tmp_image.exists():
                    tmp_image.unlink()
            except Exception:
                pass
    if page_count > MAX_PDF_OCR_PAGES:
        results.append(
            OCRResult(
                source_visual_path=str(pdf_path),
                status="skipped",
                ocr_text="",
                model_notes=f"PDF OCR limited to first {MAX_PDF_OCR_PAGES} pages; {page_count - MAX_PDF_OCR_PAGES} pages not processed",
                page_ref="remaining_pages",
            )
        )
    return results


def run_ocr_for_exports(
    export_root: Path,
    output_root: Path,
    backend: str = "local",
    display_root: Path | None = None,
) -> list[OCRResult]:
    output_root.mkdir(parents=True, exist_ok=True)
    export_root = export_root.expanduser().resolve()
    image_paths = [
        p
        for p in export_root.rglob("*")
        if p.is_file() and p.suffix.lower() in {".png", ".jpg", ".jpeg", ".bmp", ".webp"}
    ]
    pdf_paths = [p for p in export_root.rglob("*.pdf") if p.is_file()]
    results: list[OCRResult] = []
    for image in image_paths:
        result = _run_local_tesseract(image) if backend == "local" else OCRResult(str(image), "skipped", "", "unsupported backend")
        if display_root is not None:
            result.source_visual_path = _display_output_path(Path(result.source_visual_path), display_root)
        results.append(result)
        ocr_name = f"{_artifact_stem(str(image.relative_to(export_root)))}.ocr.json"
        (output_root / ocr_name).write_text(json.dumps(asdict(result), ensure_ascii=False, indent=2), encoding="utf-8")
    for pdf_path in pdf_paths:
        pdf_results = _run_pdf_ocr(pdf_path, backend)
        if display_root is not None:
            for item in pdf_results:
                item.source_visual_path = _display_output_path(Path(item.source_visual_path), display_root)
        results.extend(pdf_results)
        pdf_stem = _artifact_stem(str(pdf_path.relative_to(export_root)))
        for idx, item in enumerate(pdf_results, start=1):
            (output_root / f"{pdf_stem}.page{idx}.ocr.json").write_text(
                json.dumps(asdict(item), ensure_ascii=False, indent=2),
                encoding="utf-8",
            )
    return results


def stage_attachment_files(
    attachment_paths: list[Path],
    visual_export_root: Path,
    source_labels: list[str] | None = None,
) -> list[Path]:
    staged: list[Path] = []
    attachment_dir = visual_export_root / "attachments"
    attachment_dir.mkdir(parents=True, exist_ok=True)
    for idx, source in enumerate(attachment_paths):
        label = source_labels[idx] if source_labels and idx < len(source_labels) else source.name
        target = attachment_dir / f"{_artifact_stem(label)}{source.suffix.lower()}"
        try:
            shutil.copy2(source, target)
            staged.append(target)
        except Exception:
            continue
    return staged


def write_file_inventory(path: Path, rows: list[FileInventoryRow]) -> None:
    lines = [
        "# file_inventory",
        "",
        "| relative_path | file_name | file_type | size_bytes | modified_time | should_process | skip_reason |",
        "|---|---|---:|---:|---|---|---|",
    ]
    for row in rows:
        lines.append(
            f"| {_md_cell(row.relative_path)} | {_md_cell(row.file_name)} | {_md_cell(row.file_type)} | {row.size_bytes} | "
            f"{_md_cell(row.modified_time)} | {row.should_process} | {_md_cell(row.skip_reason or '-')} |"
        )
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def _count_attr(sheet: SheetAnalysis, attr: str) -> int:
    return sum(1 for cell in sheet.cell_records if getattr(cell, attr))


def write_workbook_inventory(path: Path, books: list[WorkbookAnalysis]) -> None:
    lines = ["# workbook_inventory", ""]
    for book in books:
        lines.extend(
            [
                f"## {book.workbook_name}",
                f"- source_file: `{book.source_file}`",
                f"- sheet_order: {', '.join(book.sheet_order) if book.sheet_order else '(none)'}",
                f"- hidden_sheets: {', '.join(book.hidden_sheets) if book.hidden_sheets else '(none)'}",
                f"- named_ranges: {len(book.named_ranges)}",
                "",
            ]
        )
        for sheet in book.sheet_analyses:
            lines.extend(
                [
                    f"### sheet: {sheet.sheet_name}",
                    f"- state: {sheet.state}",
                    f"- dimensions: {sheet.dimensions}",
                    f"- max_row/max_column: {sheet.max_row}/{sheet.max_column}",
                    f"- merged_ranges: {len(sheet.merged_ranges)}",
                    f"- data_validation_count: {sheet.data_validation_count}",
                    f"- conditional_format_count: {sheet.conditional_format_count}",
                    f"- chart_count: {sheet.chart_count}",
                    f"- image_count: {sheet.image_count}",
                    f"- hyperlinks/comments/formulas: {_count_attr(sheet, 'hyperlink')}/{_count_attr(sheet, 'comment')}/{_count_attr(sheet, 'formula')}",
                    "",
                ]
            )
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def write_document_inventory(path: Path, docs: list[DocumentAnalysis]) -> None:
    lines = ["# document_inventory", ""]
    for doc in docs:
        lines.extend(
            [
                f"## {doc.document_name}",
                f"- source_file: `{doc.source_file}`",
                f"- document_type: {doc.document_type}",
                f"- section_count: {len(doc.sections)}",
                f"- image_count: {doc.image_count}",
                f"- warnings: {len(doc.extraction_warnings)}",
                "",
            ]
        )
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def _scan_keywords_from_workbooks(books: list[WorkbookAnalysis], keywords: tuple[str, ...], cap: int) -> list[str]:
    hits: list[str] = []
    low_keys = tuple(k.lower() for k in keywords)
    for book in books:
        for sheet in book.sheet_analyses:
            for cell in sheet.cell_records[:800]:
                if not cell.value:
                    continue
                value = cell.value.lower()
                if any(k in value for k in low_keys):
                    hits.append(f"{book.workbook_name}:{sheet.sheet_name}.{cell.coordinate}:{cell.value}")
                    if len(hits) >= cap:
                        return hits
    return hits


def _scan_keywords_from_docs(docs: list[DocumentAnalysis], keywords: tuple[str, ...], cap: int) -> list[str]:
    hits: list[str] = []
    low_keys = tuple(k.lower() for k in keywords)
    for doc in docs:
        for section in doc.sections[:1200]:
            value = section.text.lower()
            if any(k in value for k in low_keys):
                hits.append(f"{doc.document_name}:{section.location}:{section.text}")
                if len(hits) >= cap:
                    return hits
    return hits


def _scan_keywords(workbooks: list[WorkbookAnalysis], docs: list[DocumentAnalysis], keywords: tuple[str, ...], cap: int) -> list[str]:
    hits = _scan_keywords_from_workbooks(workbooks, keywords, cap)
    if len(hits) < cap:
        hits.extend(_scan_keywords_from_docs(docs, keywords, cap - len(hits)))
    return hits


def write_deep_notes(output_dir: Path, books: list[WorkbookAnalysis], docs: list[DocumentAnalysis]) -> None:
    output_dir.mkdir(parents=True, exist_ok=True)
    for book in books:
        lines = [f"# deep_reading: {book.workbook_name}", "", "## Observations"]
        if not book.sheet_analyses:
            lines.append("- No readable sheet content.")
        for sheet in book.sheet_analyses:
            lines.append(f"- `{sheet.sheet_name}` contains {len(sheet.cell_records)} non-empty/comment/link cells.")
        if book.extraction_warnings:
            lines.append("")
            lines.append("## Warnings")
            lines.extend([f"- {x}" for x in book.extraction_warnings])
        (output_dir / f"{_artifact_stem(book.source_file)}.md").write_text("\n".join(lines) + "\n", encoding="utf-8")
    for doc in docs:
        lines = [
            f"# deep_reading: {doc.document_name}",
            "",
            "## Observations",
            f"- document_type: {doc.document_type}",
            f"- section_count: {len(doc.sections)}",
            f"- image_count: {doc.image_count}",
            "",
            "## Section Samples",
        ]
        if not doc.sections:
            lines.append("- No text sections extracted.")
        else:
            for section in doc.sections[:80]:
                lines.append(f"- {section.section_type} {section.location}: {section.text}")
        if doc.extraction_warnings:
            lines.append("")
            lines.append("## Warnings")
            lines.extend([f"- {x}" for x in doc.extraction_warnings])
        (output_dir / f"{_artifact_stem(doc.source_file)}.md").write_text("\n".join(lines) + "\n", encoding="utf-8")


def _purpose(workbooks: list[WorkbookAnalysis], docs: list[DocumentAnalysis]) -> str:
    flow = _scan_keywords(workbooks, docs, ("flow", "step", "手順", "処理"), 1)
    io = _scan_keywords(workbooks, docs, ("input", "output", "入力", "出力"), 1)
    if flow:
        return "业务处理手顺/流程说明"
    if io:
        return "输入输出与处理条件定义"
    return "推定：业务说明或设计资料"


def write_final_summary(path: Path, books: list[WorkbookAnalysis], docs: list[DocumentAnalysis]) -> None:
    lines = ["# final_summary", ""]
    for book in books:
        lines.extend(
            [
                f"## {book.workbook_name}",
                f"- 文件名: {book.workbook_name}",
                f"- 文件目的: {_purpose([book], [])}",
                "- 适用业务/系统: 不确定",
                f"- 主要sheet: {', '.join(book.sheet_order) if book.sheet_order else '(none)'}",
                f"- 重要流程: {'; '.join(_scan_keywords([book], [], ('flow', 'step', '手順', '処理'), 6)) or '不确定'}",
                f"- Input: {'; '.join(_scan_keywords([book], [], ('input', '入力', '検索条件', '条件', 'file', 'path'), 8)) or '不确定'}",
                f"- Output: {'; '.join(_scan_keywords([book], [], ('output', '出力', 'result', '一覧', 'download'), 8)) or '不确定'}",
                f"- 系统/画面操作: {'; '.join(_scan_keywords([book], [], ('click', 'button', 'screen', '画面', '押下', '遷移'), 8)) or '不确定'}",
                f"- 数据更新/查询/下载动作: {'; '.join(_scan_keywords([book], [], ('insert', 'update', 'delete', 'register', '更新', '削除', '照会', 'download'), 8)) or '不确定'}",
                f"- 条件分支: {'; '.join(_scan_keywords([book], [], ('if', 'else', 'when', '条件', '場合', '判定', '分岐'), 8)) or '不确定'}",
                f"- 异常处理: {'; '.join(_scan_keywords([book], [], ('error', 'exception', 'failed', '警告', 'エラー', '異常'), 8)) or '不确定'}",
                f"- 关键字段: {'; '.join(_scan_keywords([book], [], ('id', 'code', 'name', 'date', 'status', 'flag'), 10)) or '不确定'}",
                f"- 关键截图/OCR结论: {'Workbook contains embedded images.' if sum(s.image_count for s in book.sheet_analyses) else 'No embedded image evidence captured.'}",
                f"- 未确认事项: {'; '.join(book.extraction_warnings[:6]) if book.extraction_warnings else 'None observed at current extraction depth.'}",
                f"- 可信度: {'high' if sum(len(s.cell_records) for s in book.sheet_analyses) > 200 else ('medium' if sum(len(s.cell_records) for s in book.sheet_analyses) > 50 else 'low')}",
                "",
            ]
        )
    for doc in docs:
        lines.extend(
            [
                f"## {doc.document_name}",
                f"- 文件名: {doc.document_name}",
                f"- 文件目的: {_purpose([], [doc])}",
                "- 适用业务/系统: 不确定",
                f"- 主要sheet: (document:{doc.document_type})",
                f"- 重要流程: {'; '.join(_scan_keywords([], [doc], ('flow', 'step', '手順', '処理'), 6)) or '不确定'}",
                f"- Input: {'; '.join(_scan_keywords([], [doc], ('input', '入力', '検索条件', '条件', 'file', 'path'), 8)) or '不确定'}",
                f"- Output: {'; '.join(_scan_keywords([], [doc], ('output', '出力', 'result', '一覧', 'download'), 8)) or '不确定'}",
                f"- 系统/画面操作: {'; '.join(_scan_keywords([], [doc], ('click', 'button', 'screen', '画面', '押下', '遷移'), 8)) or '不确定'}",
                f"- 数据更新/查询/下载动作: {'; '.join(_scan_keywords([], [doc], ('insert', 'update', 'delete', 'register', '更新', '削除', '照会', 'download'), 8)) or '不确定'}",
                f"- 条件分支: {'; '.join(_scan_keywords([], [doc], ('if', 'else', 'when', '条件', '場合', '判定', '分岐'), 8)) or '不确定'}",
                f"- 异常处理: {'; '.join(_scan_keywords([], [doc], ('error', 'exception', 'failed', '警告', 'エラー', '異常'), 8)) or '不确定'}",
                f"- 关键字段: {'; '.join(_scan_keywords([], [doc], ('id', 'code', 'name', 'date', 'status', 'flag'), 10)) or '不确定'}",
                f"- 关键截图/OCR结论: {'Document has images.' if doc.image_count else 'No direct image evidence captured.'}",
                f"- 未确认事项: {'; '.join(doc.extraction_warnings[:6]) if doc.extraction_warnings else 'None observed at current extraction depth.'}",
                f"- 可信度: {'high' if len(doc.sections) > 200 else ('medium' if len(doc.sections) > 60 else 'low')}",
                "",
            ]
        )
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def write_structured_json(path: Path, result: PipelineResult) -> None:
    path.write_text(json.dumps(asdict(result), ensure_ascii=False, indent=2), encoding="utf-8")


def run_pipeline(config: PipelineConfig) -> PipelineResult:
    config.normalize()
    config.ensure_output_dirs()
    result = PipelineResult.bootstrap(config.input_path, config.output_root)
    result.warnings.extend(config.validation_warnings)
    inventory = scan_input(config)
    result.file_inventory = inventory
    write_file_inventory(config.output_root / "file_inventory.md", inventory)

    spreadsheet_rows = filter_spreadsheet_targets(inventory)
    document_rows = filter_document_targets(inventory)
    attachment_rows = filter_attachment_targets(inventory)

    if attachment_rows:
        attachment_paths = [
            (config.input_path / row.relative_path) if config.input_path.is_dir() else config.input_path
            for row in attachment_rows
        ]
        staged = stage_attachment_files(
            attachment_paths,
            config.visual_exports_path,
            [row.relative_path for row in attachment_rows],
        )
        result.warnings.append(f"Attachment candidates detected: {len(attachment_rows)}")
        result.warnings.append(f"Attachment staged for OCR/vision: {len(staged)}")

    def resolve_source(row: FileInventoryRow) -> Path:
        return config.input_path / row.relative_path if config.input_path.is_dir() else config.input_path

    def artifact_stem(row: FileInventoryRow) -> str:
        return _artifact_stem(row.relative_path)

    for row in spreadsheet_rows + document_rows:
        source = resolve_source(row)
        if config.enable_markitdown:
            md_target = config.extracted_markdown_path / f"{artifact_stem(row)}.md"
            result.markdown_results.append(
                extract_to_markdown(
                    source,
                    md_target,
                    source_label=row.relative_path,
                    output_label=_display_output_path(md_target, config.output_root),
                )
            )

    for row in spreadsheet_rows:
        source = resolve_source(row)
        workbook = analyze_workbook(source, row.relative_path)
        result.workbook_results.append(workbook)
        if config.enable_visual_export:
            visual_dir = config.visual_exports_path / artifact_stem(row)
            exported = export_visual_assets(source, workbook, visual_dir)
            for obj in exported:
                if obj.export_path:
                    obj.export_path = _display_output_path(Path(obj.export_path), config.output_root)
            for sheet in workbook.sheet_analyses:
                sheet.object_records.extend([obj for obj in exported if obj.sheet_name in {sheet.sheet_name, "*"}])

    for row in document_rows:
        source = resolve_source(row)
        doc = analyze_document(source, row.relative_path)
        result.document_results.append(doc)
        if config.enable_visual_export:
            export_document_visuals(source, config.visual_exports_path / artifact_stem(row))

    if result.markdown_results:
        lines = ["# markitdown_extraction_log", ""]
        for item in result.markdown_results:
            lines.extend(
                [
                    f"## {Path(item.source_file).name}",
                    f"- status: {item.status}",
                    f"- markdown_path: {item.markdown_path or '(none)'}",
                    f"- log: {item.log or '(empty)'}",
                    "",
                ]
            )
        (config.extracted_markdown_path / "extraction_log.md").write_text("\n".join(lines) + "\n", encoding="utf-8")

    if config.enable_ocr:
        result.ocr_results = run_ocr_for_exports(
            config.visual_exports_path,
            config.ocr_results_path,
            config.ocr_backend,
            display_root=config.output_root,
        )

    write_workbook_inventory(config.output_root / "workbook_inventory.md", result.workbook_results)
    write_document_inventory(config.output_root / "document_inventory.md", result.document_results)
    write_deep_notes(config.deep_notes_path, result.workbook_results, result.document_results)
    write_final_summary(config.output_root / "final_summary.md", result.workbook_results, result.document_results)
    write_structured_json(config.output_root / "structured_data.json", result)
    return result
